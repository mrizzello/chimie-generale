#!/usr/bin/env python3
"""Build a minimalist .pptx deck from a slides/content/<slug>.yml plan.

Every text-bearing slide is built on a real title/content placeholder (from
python-pptx's built-in slide layouts), never a freeform textbox, and no
font/size/color is set on placeholder runs. This is what lets the deck be
opened in Keynote/PowerPoint and re-themed: a theme change restyles
placeholders by role (title, body, subtitle) and paragraph level, which is
exactly the "hierarchy" a freeform textbox can't participate in.

Usage: build_deck.py slides/content/introduction.yml
"""
import sys
from pathlib import Path

import yaml
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

from latex_render import render_tex

REPO_ROOT = Path(__file__).resolve().parents[2]

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.9)
GAP = Inches(0.25)
RENDER_DPI = 400

LAYOUT_TITLE = 0
LAYOUT_CONTENT = 1
LAYOUT_SECTION = 2

DROP_PLACEHOLDER_IDX = {10, 11, 12}  # date, footer, slide number


def rescale_layout_placeholders(layout, scale_x, scale_y):
    for ph in layout.placeholders:
        ph.left = int(ph.left * scale_x)
        ph.width = int(ph.width * scale_x)
        ph.top = int(ph.top * scale_y)
        ph.height = int(ph.height * scale_y)


def get_placeholder(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def remove_shape(shape):
    shape._element.getparent().remove(shape._element)


def strip_footer_placeholders(slide):
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx in DROP_PLACEHOLDER_IDX:
            remove_shape(ph)


def set_bullets(placeholder, items):
    tf = placeholder.text_frame
    for i, item in enumerate(items):
        if isinstance(item, dict):
            text, level = item.get("text", ""), int(item.get("level", 0))
        else:
            text, level = str(item), 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level


def render_snippet(item, cache_dir):
    kind = item.get("kind", "math")
    if kind == "image":
        return REPO_ROOT / item["path"]
    if kind == "table":
        tex = item["tex"]
    else:
        raw = item["tex"].strip()
        if raw.startswith("\\begin{") or raw.startswith("\\["):
            tex = raw if raw.startswith("\\[") else f"\\[{raw}\\]"
        else:
            tex = f"$\\displaystyle {raw}$"
    return render_tex(tex, kind, cache_dir, dpi=RENDER_DPI)


def add_image_fit(slide, img_path, top, max_width, max_height, allow_upscale=False):
    with Image.open(img_path) as im:
        px_w, px_h = im.size
    if allow_upscale:
        # Chapter illustrations are arbitrary web-resolution rasters, not
        # DPI-calibrated renders: fit them to the box by aspect ratio alone
        # (upscaling as needed) instead of capping at a native px/DPI size.
        aspect = px_w / px_h
        box_aspect = (max_width / Inches(1)) / (max_height / Inches(1))
        if aspect > box_aspect:
            w, h = max_width, Inches((max_width / Inches(1)) / aspect)
        else:
            h, w = max_height, Inches((max_height / Inches(1)) * aspect)
    else:
        w_in = px_w / RENDER_DPI
        h_in = px_h / RENDER_DPI
        scale = min((max_width / Inches(1)) / w_in, (max_height / Inches(1)) / h_in, 1.0)
        w, h = Inches(w_in * scale), Inches(h_in * scale)
    left = MARGIN + (max_width - w) / 2
    slide.shapes.add_picture(str(img_path), left, top, width=w, height=h)
    return h


def place_renders(slide, renders, cache_dir, top, bottom_limit):
    if not renders:
        return
    remaining_h = bottom_limit - top
    per_slot = remaining_h / len(renders)
    for item in renders:
        png = render_snippet(item, cache_dir)
        allow_upscale = item.get("kind") == "image"
        add_image_fit(slide, png, top, SLIDE_W - 2 * MARGIN, per_slot - GAP, allow_upscale)
        top += per_slot


def add_text_slide(prs, chapter_num, exercise_counter, slide_spec, cache_dir):
    stype = slide_spec["type"]
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    strip_footer_placeholders(slide)

    if stype == "exercise":
        exercise_counter += 1
        title_text = f"Exercice {chapter_num}.{exercise_counter}"
    elif stype == "answer":
        title_text = f"Corrigé {chapter_num}.{exercise_counter}"
    elif stype == "objectives":
        title_text = "Objectifs"
    else:
        title_text = slide_spec.get("title") or stype.capitalize()

    label = slide_spec.get("label")
    if label:
        title_text = f"{title_text} — {label}"
    slide.shapes.title.text = title_text

    body_ph = get_placeholder(slide, 1)
    items = slide_spec.get("items")
    renders = slide_spec.get("renders")
    content_bottom_limit = SLIDE_H - Inches(0.4)

    if items:
        set_bullets(body_ph, items)
        if renders:
            body_ph.height = Inches(min(
                body_ph.height / Inches(1),
                (content_bottom_limit - body_ph.top) * 0.42 / Inches(1),
            ))
        image_top = body_ph.top + body_ph.height + GAP
    else:
        image_top = body_ph.top if body_ph is not None else Inches(1.75)
        if body_ph is not None:
            remove_shape(body_ph)

    if renders:
        place_renders(slide, renders, cache_dir, image_top, content_bottom_limit)

    return exercise_counter


def build_deck(content: dict, out_path: Path, cache_dir: Path):
    prs = Presentation()
    scale_x = SLIDE_W / prs.slide_width
    scale_y = SLIDE_H / prs.slide_height
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for idx in (LAYOUT_TITLE, LAYOUT_CONTENT, LAYOUT_SECTION):
        rescale_layout_placeholders(prs.slide_layouts[idx], scale_x, scale_y)

    chapter_num = int(content.get("number", "0") or 0) // 10
    exercise_counter = 0

    for slide_spec in content["slides"]:
        stype = slide_spec["type"]

        if stype == "title":
            slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
            strip_footer_placeholders(slide)
            slide.shapes.title.text = slide_spec.get("title") or content["chapter"]
            subtitle = get_placeholder(slide, 1)
            if subtitle is not None:
                subtitle.text = f"Chapitre {chapter_num}"

        elif stype == "divider":
            slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_SECTION])
            strip_footer_placeholders(slide)
            slide.shapes.title.text = slide_spec["title"]
            body = get_placeholder(slide, 1)
            if body is not None:
                remove_shape(body)

        else:
            exercise_counter = add_text_slide(
                prs, chapter_num, exercise_counter, slide_spec, cache_dir
            )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("usage: build_deck.py <content.yml>", file=sys.stderr)
        sys.exit(1)
    yml_path = Path(sys.argv[1])
    content = yaml.safe_load(yml_path.read_text())
    slug = content["slug"]
    number = content.get("number", "")
    out_path = REPO_ROOT / "slides" / f"{number}-{slug}.pptx"
    cache_dir = REPO_ROOT / "slides" / "_assets" / slug
    build_deck(content, out_path, cache_dir)
    print(f"Wrote {out_path}")
